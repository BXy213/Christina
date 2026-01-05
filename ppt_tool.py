"""
PPT 生成工具
根据用户提纲自动生成 PowerPoint 演示文稿
"""
import os
import re
import json
import requests
import tempfile
from io import BytesIO
from pathlib import Path
from datetime import datetime
from typing import Optional, List, Dict, Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from langchain_core.tools import BaseTool
from langchain_openai import ChatOpenAI
from pydantic import BaseModel, Field

from config import config
from logger import logger


# ============== 输入参数模型 ==============

class PPTInput(BaseModel):
    """PPT 生成工具的输入参数"""
    outline: str = Field(
        description="PPT 的主题和提纲内容。可以是自然语言描述，包含要点、章节标题等。"
                    "支持图片引用格式：[图片: 本地路径] 或 [图片: URL]"
    )
    output_name: Optional[str] = Field(
        default=None,
        description="输出文件名（不含扩展名），默认使用主题名+时间戳"
    )


# ============== PPT 样式主题 ==============

class PPTTheme:
    """PPT 样式主题配置"""
    
    def __init__(self):
        primary = config.PPT_PRIMARY_COLOR
        accent = config.PPT_ACCENT_COLOR
        
        self.primary = RGBColor(primary[0], primary[1], primary[2])
        self.accent = RGBColor(accent[0], accent[1], accent[2])
        self.text_dark = RGBColor(51, 51, 51)
        self.text_light = RGBColor(255, 255, 255)
        self.background = RGBColor(255, 255, 255)
        self.font_name = config.PPT_FONT_NAME
        self.title_size = Pt(config.PPT_TITLE_FONT_SIZE)
        self.content_size = Pt(config.PPT_CONTENT_FONT_SIZE)
        self.subtitle_size = Pt(28)
        self.bullet_size = Pt(20)


# ============== LLM 提纲解析 Prompt ==============

OUTLINE_PARSER_PROMPT = """你是一个专业的 PPT 制作助手。请将用户的提纲解析为结构化 JSON。

## 你的任务
将用户提供的自然语言提纲转换为 PPT 页面结构，包括：
1. 识别主题和各章节
2. 智能分页（语义相关的内容放一页，每页2-5个要点）
3. 为每页选择合适的布局
4. 识别图片引用并确定其位置

## 分页规则
- 主题明显切换时分页
- 每页内容不超过5个要点
- 有配图的内容适合单独成页
- "感谢"、"总结"、"结束"等词识别为结尾页

## 图片引用格式识别
- [图片: D:\\path\\to\\image.png] → type: "local", path: "D:\\path\\to\\image.png"
- [图片: https://example.com/img.jpg] → type: "url", url: "https://example.com/img.jpg"
- [图片: C:/Users/xxx/photo.jpg] → type: "local", path: "C:/Users/xxx/photo.jpg"

## 布局类型
- title_slide: 封面页（有主标题和副标题）
- bullet_list: 纯文字列表页
- image_right: 左文右图
- image_left: 左图右文
- image_full: 大图配标题
- two_column: 双栏对比
- ending_slide: 结尾页

## 输出格式（严格 JSON）
```json
{
  "title": "PPT主标题",
  "slides": [
    {
      "type": "title_slide",
      "title": "主标题",
      "subtitle": "副标题（可选）"
    },
    {
      "type": "bullet_list",
      "title": "页面标题",
      "content": ["要点1", "要点2", "要点3"]
    },
    {
      "type": "image_right",
      "title": "页面标题",
      "content": ["要点1", "要点2"],
      "image": {
        "type": "local",
        "path": "D:\\\\images\\\\example.png"
      }
    },
    {
      "type": "two_column",
      "title": "对比标题",
      "left_title": "左栏标题",
      "left_content": ["左1", "左2"],
      "right_title": "右栏标题",
      "right_content": ["右1", "右2"]
    },
    {
      "type": "ending_slide",
      "title": "感谢观看",
      "subtitle": "联系方式（可选）"
    }
  ]
}
```

## 重要提示
- 只输出 JSON，不要有其他文字
- 确保 JSON 格式正确
- Windows 路径中的反斜杠需要转义为 \\\\
- 如果用户没有明确分页，根据语义自动分组

用户提纲：
{outline}
"""


# ============== PPT 生成工具 ==============

class PPTGeneratorTool(BaseTool):
    """PPT 生成工具"""
    
    name: str = "ppt_generator"
    description: str = """
    根据用户提供的提纲生成 PPT 演示文稿。
    
    功能：
    - 支持自然语言提纲，自动分页和排版
    - 支持图片引用：[图片: 本地路径] 或 [图片: URL]
    - 自动选择合适的页面布局
    - 生成专业的 PPTX 文件
    
    使用场景：
    - 用户说"帮我做一个关于XX的PPT"
    - 用户提供了具体的提纲内容
    - 用户需要快速生成演示文稿
    """
    args_schema: type[BaseModel] = PPTInput
    
    # 类属性声明
    llm: Any = None
    theme: PPTTheme = None
    output_dir: Path = None
    
    def __init__(self, llm: ChatOpenAI, **kwargs):
        """
        初始化 PPT 生成工具
        
        Args:
            llm: LangChain LLM 实例，用于解析提纲
        """
        super().__init__(**kwargs)
        self.llm = llm
        self.theme = PPTTheme()
        self.output_dir = Path(config.PPT_OUTPUT_DIR)
        
        # 确保输出目录存在
        self.output_dir.mkdir(parents=True, exist_ok=True)
    
    def _run(self, outline: str, output_name: Optional[str] = None) -> str:
        """
        执行 PPT 生成
        
        Args:
            outline: 用户提供的提纲
            output_name: 输出文件名
            
        Returns:
            生成结果信息
        """
        try:
            logger.log("正在解析提纲...")
            
            # 1. 使用 LLM 解析提纲
            structured_data = self._parse_outline_with_llm(outline)
            if not structured_data:
                return "解析提纲失败，请检查提纲格式是否正确。"
            
            logger.log(f"解析完成，共 {len(structured_data.get('slides', []))} 页")
            
            # 2. 处理图片
            structured_data = self._process_images(structured_data)
            
            # 3. 渲染 PPT
            pptx_path = self._render_pptx(structured_data, output_name)
            
            # 4. 返回结果
            slide_count = len(structured_data.get('slides', []))
            return f"✅ PPT 已成功生成！\n📄 文件路径：{pptx_path}\n📊 共 {slide_count} 页"
            
        except Exception as e:
            logger.error(f"PPT 生成失败: {e}")
            return f"PPT 生成失败：{str(e)}"
    
    async def _arun(self, outline: str, output_name: Optional[str] = None) -> str:
        """异步执行（使用同步实现）"""
        return self._run(outline, output_name)
    
    def _parse_outline_with_llm(self, outline: str) -> Optional[Dict]:
        """
        使用 LLM 解析自然语言提纲为结构化数据
        
        Args:
            outline: 用户提纲
            
        Returns:
            结构化的 PPT 数据
        """
        try:
            prompt = OUTLINE_PARSER_PROMPT.format(outline=outline)
            response = self.llm.invoke(prompt)
            
            # 提取 JSON
            content = response.content
            
            # 尝试从 markdown 代码块中提取
            json_match = re.search(r'```(?:json)?\s*([\s\S]*?)\s*```', content)
            if json_match:
                json_str = json_match.group(1)
            else:
                json_str = content
            
            # 解析 JSON
            data = json.loads(json_str)
            return data
            
        except json.JSONDecodeError as e:
            logger.error(f"JSON 解析失败: {e}")
            logger.error(f"LLM 返回内容: {content[:500]}")
            return None
        except Exception as e:
            logger.error(f"LLM 调用失败: {e}")
            return None
    
    def _process_images(self, data: Dict) -> Dict:
        """
        处理所有图片引用
        
        Args:
            data: 结构化 PPT 数据
            
        Returns:
            处理后的数据（图片转为本地路径）
        """
        for slide in data.get('slides', []):
            if 'image' not in slide:
                continue
                
            img = slide['image']
            img_type = img.get('type', '')
            
            if img_type == 'url':
                # 下载 URL 图片
                url = img.get('url', '')
                local_path = self._download_image(url)
                if local_path:
                    img['local_path'] = local_path
                else:
                    logger.warning(f"图片下载失败: {url}")
                    
            elif img_type == 'local':
                # 验证本地路径
                path = img.get('path', '')
                if os.path.exists(path):
                    img['local_path'] = path
                else:
                    logger.warning(f"本地图片不存在: {path}")
        
        return data
    
    def _download_image(self, url: str) -> Optional[str]:
        """
        下载网络图片到临时文件
        
        Args:
            url: 图片 URL
            
        Returns:
            本地临时文件路径
        """
        try:
            response = requests.get(url, timeout=10)
            response.raise_for_status()
            
            # 获取文件扩展名
            content_type = response.headers.get('content-type', '')
            if 'png' in content_type:
                ext = '.png'
            elif 'gif' in content_type:
                ext = '.gif'
            else:
                ext = '.jpg'
            
            # 保存到临时文件
            temp_file = tempfile.NamedTemporaryFile(
                delete=False, 
                suffix=ext,
                dir=str(self.output_dir)
            )
            temp_file.write(response.content)
            temp_file.close()
            
            return temp_file.name
            
        except Exception as e:
            logger.error(f"下载图片失败 {url}: {e}")
            return None
    
    def _render_pptx(self, data: Dict, output_name: Optional[str]) -> str:
        """
        使用 python-pptx 渲染 PPT
        
        Args:
            data: 结构化 PPT 数据
            output_name: 输出文件名
            
        Returns:
            生成的文件路径
        """
        prs = Presentation()
        
        # 设置幻灯片尺寸 (16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # 渲染每一页
        for slide_data in data.get('slides', []):
            slide_type = slide_data.get('type', 'bullet_list')
            
            if slide_type == 'title_slide':
                self._add_title_slide(prs, slide_data)
            elif slide_type == 'bullet_list':
                self._add_bullet_slide(prs, slide_data)
            elif slide_type == 'image_right':
                self._add_image_slide(prs, slide_data, image_position='right')
            elif slide_type == 'image_left':
                self._add_image_slide(prs, slide_data, image_position='left')
            elif slide_type == 'image_full':
                self._add_full_image_slide(prs, slide_data)
            elif slide_type == 'two_column':
                self._add_two_column_slide(prs, slide_data)
            elif slide_type == 'ending_slide':
                self._add_ending_slide(prs, slide_data)
            else:
                # 默认使用列表布局
                self._add_bullet_slide(prs, slide_data)
        
        # 生成文件名
        if not output_name:
            title = data.get('title', 'presentation')
            # 清理文件名中的非法字符
            title = re.sub(r'[<>:"/\\|?*]', '', title)
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            output_name = f"{title}_{timestamp}"
        
        # 保存文件
        output_path = self.output_dir / f"{output_name}.pptx"
        prs.save(str(output_path))
        
        return str(output_path)
    
    def _add_title_slide(self, prs: Presentation, data: Dict):
        """添加封面页"""
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加背景色块
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.theme.primary
        shape.line.fill.background()
        
        # 添加装饰条
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(5),
            prs.slide_width, Inches(0.1)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self.theme.accent
        accent_bar.line.fill.background()
        
        # 添加标题
        title = data.get('title', '演示文稿')
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5),
            Inches(11.333), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(54)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = self.theme.text_light
        tf.paragraphs[0].font.name = self.theme.font_name
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 添加副标题
        subtitle = data.get('subtitle', '')
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.2),
                Inches(11.333), Inches(0.8)
            )
            tf = subtitle_box.text_frame
            tf.paragraphs[0].text = subtitle
            tf.paragraphs[0].font.size = self.theme.subtitle_size
            tf.paragraphs[0].font.color.rgb = self.theme.text_light
            tf.paragraphs[0].font.name = self.theme.font_name
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _add_bullet_slide(self, prs: Presentation, data: Dict):
        """添加列表页"""
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加顶部色带
        header_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, Inches(1.2)
        )
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = self.theme.primary
        header_bar.line.fill.background()
        
        # 添加标题
        title = data.get('title', '')
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = self.theme.title_size
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = self.theme.text_light
        tf.paragraphs[0].font.name = self.theme.font_name
        
        # 添加内容要点
        content = data.get('content', [])
        if content:
            content_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(1.8),
                Inches(11.5), Inches(5)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            
            for i, item in enumerate(content):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = f"● {item}"
                p.font.size = self.theme.content_size
                p.font.color.rgb = self.theme.text_dark
                p.font.name = self.theme.font_name
                p.space_after = Pt(16)
    
    def _add_image_slide(self, prs: Presentation, data: Dict, image_position: str = 'right'):
        """添加图文混排页"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加顶部色带
        header_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, Inches(1.2)
        )
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = self.theme.primary
        header_bar.line.fill.background()
        
        # 添加标题
        title = data.get('title', '')
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = self.theme.title_size
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = self.theme.text_light
        tf.paragraphs[0].font.name = self.theme.font_name
        
        # 根据图片位置确定布局
        if image_position == 'right':
            text_left = Inches(0.5)
            text_width = Inches(6)
            img_left = Inches(7)
        else:
            text_left = Inches(6.5)
            text_width = Inches(6)
            img_left = Inches(0.5)
        
        # 添加文字内容
        content = data.get('content', [])
        if content:
            content_box = slide.shapes.add_textbox(
                text_left, Inches(1.8),
                text_width, Inches(5)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            
            for i, item in enumerate(content):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = f"● {item}"
                p.font.size = self.theme.content_size
                p.font.color.rgb = self.theme.text_dark
                p.font.name = self.theme.font_name
                p.space_after = Pt(12)
        
        # 添加图片
        img_info = data.get('image', {})
        local_path = img_info.get('local_path')
        if local_path and os.path.exists(local_path):
            try:
                slide.shapes.add_picture(
                    local_path,
                    img_left, Inches(1.8),
                    width=Inches(5.5)
                )
            except Exception as e:
                logger.warning(f"插入图片失败: {e}")
                # 添加占位符
                self._add_image_placeholder(slide, img_left, Inches(1.8), Inches(5.5), Inches(4))
        else:
            # 添加占位符
            self._add_image_placeholder(slide, img_left, Inches(1.8), Inches(5.5), Inches(4))
    
    def _add_full_image_slide(self, prs: Presentation, data: Dict):
        """添加大图页"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        title = data.get('title', '')
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = self.theme.title_size
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = self.theme.primary
        tf.paragraphs[0].font.name = self.theme.font_name
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 添加图片
        img_info = data.get('image', {})
        local_path = img_info.get('local_path')
        
        if local_path and os.path.exists(local_path):
            try:
                # 居中放置大图
                slide.shapes.add_picture(
                    local_path,
                    Inches(1.5), Inches(1.5),
                    width=Inches(10)
                )
            except Exception as e:
                logger.warning(f"插入图片失败: {e}")
                self._add_image_placeholder(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(5))
        else:
            self._add_image_placeholder(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(5))
    
    def _add_two_column_slide(self, prs: Presentation, data: Dict):
        """添加双栏对比页"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加顶部色带
        header_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, Inches(1.2)
        )
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = self.theme.primary
        header_bar.line.fill.background()
        
        # 添加主标题
        title = data.get('title', '')
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = self.theme.title_size
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = self.theme.text_light
        tf.paragraphs[0].font.name = self.theme.font_name
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 左栏
        left_title = data.get('left_title', '左栏')
        left_content = data.get('left_content', [])
        
        # 左栏标题
        left_title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(5.8), Inches(0.6)
        )
        tf = left_title_box.text_frame
        tf.paragraphs[0].text = left_title
        tf.paragraphs[0].font.size = self.theme.subtitle_size
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = self.theme.primary
        tf.paragraphs[0].font.name = self.theme.font_name
        
        # 左栏内容
        left_content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.2),
            Inches(5.8), Inches(4.5)
        )
        tf = left_content_box.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(left_content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"● {item}"
            p.font.size = self.theme.bullet_size
            p.font.color.rgb = self.theme.text_dark
            p.font.name = self.theme.font_name
            p.space_after = Pt(10)
        
        # 分隔线
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(6.5), Inches(1.5),
            Inches(0.05), Inches(5)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = self.theme.accent
        divider.line.fill.background()
        
        # 右栏
        right_title = data.get('right_title', '右栏')
        right_content = data.get('right_content', [])
        
        # 右栏标题
        right_title_box = slide.shapes.add_textbox(
            Inches(7), Inches(1.5),
            Inches(5.8), Inches(0.6)
        )
        tf = right_title_box.text_frame
        tf.paragraphs[0].text = right_title
        tf.paragraphs[0].font.size = self.theme.subtitle_size
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = self.theme.primary
        tf.paragraphs[0].font.name = self.theme.font_name
        
        # 右栏内容
        right_content_box = slide.shapes.add_textbox(
            Inches(7), Inches(2.2),
            Inches(5.8), Inches(4.5)
        )
        tf = right_content_box.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(right_content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"● {item}"
            p.font.size = self.theme.bullet_size
            p.font.color.rgb = self.theme.text_dark
            p.font.name = self.theme.font_name
            p.space_after = Pt(10)
    
    def _add_ending_slide(self, prs: Presentation, data: Dict):
        """添加结尾页"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加背景
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.theme.primary
        shape.line.fill.background()
        
        # 添加装饰条
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(5), Inches(4.5),
            Inches(3.333), Inches(0.08)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self.theme.accent
        accent_bar.line.fill.background()
        
        # 添加标题
        title = data.get('title', '感谢观看')
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.8),
            Inches(11.333), Inches(1.2)
        )
        tf = title_box.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(60)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = self.theme.text_light
        tf.paragraphs[0].font.name = self.theme.font_name
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 添加副标题
        subtitle = data.get('subtitle', '')
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(5),
                Inches(11.333), Inches(0.8)
            )
            tf = subtitle_box.text_frame
            tf.paragraphs[0].text = subtitle
            tf.paragraphs[0].font.size = self.theme.content_size
            tf.paragraphs[0].font.color.rgb = self.theme.text_light
            tf.paragraphs[0].font.name = self.theme.font_name
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _add_image_placeholder(self, slide, left, top, width, height):
        """添加图片占位符"""
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(240, 240, 240)
        placeholder.line.color.rgb = RGBColor(200, 200, 200)
        
        # 添加占位文字
        text_box = slide.shapes.add_textbox(
            left, top + height / 2 - Inches(0.2),
            width, Inches(0.4)
        )
        tf = text_box.text_frame
        tf.paragraphs[0].text = "[图片位置]"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER


# ============== 工具创建函数 ==============

def create_ppt_tool(llm: ChatOpenAI) -> PPTGeneratorTool:
    """
    创建 PPT 生成工具实例
    
    Args:
        llm: LangChain LLM 实例
        
    Returns:
        PPTGeneratorTool 实例
    """
    return PPTGeneratorTool(llm=llm)


# ============== 测试代码 ==============

if __name__ == "__main__":
    from langchain_openai import ChatOpenAI
    
    # 初始化 LLM
    llm = ChatOpenAI(
        model=config.MODEL_NAME,
        temperature=config.TEMPERATURE,
        api_key=config.OPENAI_API_KEY
    )
    
    # 创建工具
    tool = create_ppt_tool(llm)
    
    # 测试提纲
    test_outline = """
    主题：人工智能发展概述
    
    介绍AI的定义和历史
    机器学习和深度学习的区别
    
    当前热门应用：
    - ChatGPT 对话系统
    - 图像识别技术
    - 自动驾驶
    
    未来展望和挑战
    
    感谢观看
    """
    
    # 生成 PPT
    result = tool._run(test_outline, "test_ai_ppt")
    logger.log(result)

